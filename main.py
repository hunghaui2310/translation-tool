# pip install python-pptx
from typing import List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET
import re
# import char
import glob
# pip install requests
import requests
import validators
import re
import docx
from python_pptx_text_replacer import TextReplacer

def call_api(text_translate):
    print(text_translate)
    # return text_translate
    query = {'sl': 'auto', 'tl': 'vi', 'dt': 't', 'client': 'gtx', 'q': text_translate}
    response = requests.get("https://translate.googleapis.com/translate_a/single", query)
    if response.ok:
        parse_response = response.json()
        last_index = len(parse_response) - 1
        if parse_response[last_index][0] == 'vi':
            return ""
        return parse_response[0][0][0]
    else:
        return ""

def translate_text_in_docx(file_path):
    doc = docx.Document(file_path)
    for p in doc.paragraphs:
        # Replace the old text with the new text
        if is_need_translate(p.text):
            p.text = call_api(p.text)

    # Save the updated document
    doc.save('output.docx')

def read_shape(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if is_need_translate(run.text):
                            pass
                            # run.text = call_api(run.text)
                            # save_text_to_file(run.text)

        # Check if the shape is a GraphicalFrame
        if shape.shape_type == 7:
            graphical_frame = shape._element
            frame_xml = graphical_frame.xml
            # Parse the XML data to extract the text
            root = ET.fromstring(frame_xml)
            for elem in root.iter():
                if is_need_translate(elem.text):
                    # elem.text = call_api(elem.text)
                    # print(elem.text)
                    pass

        # Check if the shape is a group
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            read_shape(shape.shapes)
        # Check if the shape is a table
        if shape.has_table:
            tbl = shape.table
            row_count = len(tbl.rows)
            col_count = len(tbl.columns)
            for r in range(0, row_count):
                for c in range(0, col_count):
                    cell = tbl.cell(r, c)
                    paragraphs = cell.text_frame.paragraphs
                    for paragraph in paragraphs:
                        if is_need_translate(paragraph.text):
                            save_text_to_file('test.txt', paragraph.text)
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if is_need_translate(paragraph.text):
                save_text_to_file('test.txt', paragraph.text)
                # font = paragraph.font
                # paragraph.text = call_api(paragraph.text)
                # paragraph.font.color = font.color
                # for run in paragraph.runs:
                #     if is_need_translate(run.text):
                #         run.text = call_api(run.text)

def save_text_to_file(file, text):
    with open(file, "a") as myfile:
        myfile.write(text + "\n")

def read_file_by_line(file_path):
    with open(file_path, "r") as myfile:
        data = myfile.read().splitlines()
        replacer = TextReplacer('input.pptx', slides='', tables=True, charts=True, textframes=True)
        arr = []
        for text in data:
            text_out = call_api(text)
            if text != text_out:
                # save_text_to_file('translated.txt', text_out)
                # replace_text(text, text_out)
                arr = [(text.strip(), text_out.strip()), *arr]
        replacer.replace_text(arr)
        replacer.write_presentation_to_file('./changed.pptx')

def is_need_translate(origin_text):
    if any(char.isalpha() for char in origin_text):
        if is_not_link(origin_text):
            return True
    return False


def is_not_link(origin_text):
    regex = re.compile(
        r'^(?:http|ftp)s?://'  # http:// or https://
        r'(?:(?:[A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+(?:[A-Z]{2,6}\.?|[A-Z0-9-]{2,}\.?)|'  # domain...
        r'localhost|'  # localhost...
        r'\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})'  # ...or ip
        r'(?::\d+)?'  # optional port
        r'(?:/?|[/?]\S+)$', re.IGNORECASE)
    return re.match(regex, origin_text) is None

def replace_text(old_text, new_text):
    replacer = TextReplacer('input.pptx', slides='', tables=True, charts=True, textframes=True)
    replacer.replace_text([('+ Test classification model with an input document: return True/False and list of evidences.', '+ Kiểm tra mô hình phân loại với một tài liệu đầu vào: trả về True/False và danh sách các bằng chứng.'),
                           ('+ Manage classifiers: maintain list of classifiers (Cyberbullying, Self-harm)', 'Quản lý các bộ phân loại: duy trì danh sách các bộ phân loại (Bắt nạt trên mạng, Tự làm hại bản thân)'),
                           ('+ Crawl 10.000 documents (Cyber bullying: 3.000; Self harm: 3.000; and Neutral)', '+ Thu thập thông tin 10.000 tài liệu (Bắt nạt trên mạng: 3.000; Tự hại: 3.000; và Trung lập)')])
    replacer.write_presentation_to_file('./changed2.pptx')

def search_and_replace(file_input):
    """"search and replace text in PowerPoint while preserving formatting"""
    file_output = 'output.pptx'

    prs = Presentation(file_input)
    slides = [slide for slide in prs.slides]
    shapes = []

    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)

    for slide in prs.slides:
        read_shape(slide.shapes)

    prs.save(file_output)

if __name__ == '__main__':
    file = 'test.txt'
    open(file, 'w+')
    search_and_replace('input.pptx')
    read_file_by_line(file)
    # replace_text('Project', 'Test')
